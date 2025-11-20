# app.py
import streamlit as st
from pptx import Presentation
from io import BytesIO
import json
import time

# Transformers imports are delayed until needed (to speed cold-start when using Gemini)
from typing import List, Dict, Any

st.set_page_config(page_title="Smart Presentation Builder (Free + Optional Gemini)", layout="centered")

# ---------------------------- UI: Header ---------------------------------
st.title("🎯 Smart Presentation Builder")
st.write(
    "Generate PowerPoint slides from a topic. Use built-in free HuggingFace models (no API key), "
    "or optionally use Google Gemini by pasting an API key."
)

# ---------------------------- Sidebar / Options ---------------------------
st.sidebar.header("Model & Options")
model_choice = st.sidebar.radio(
    "Choose a model (HF = Hugging Face, offline/free)",
    ("HF — Fast (flan-t5-base)", "HF — Better (flan-t5-large)", "Gemini (requires API key)")
)

num_slides = st.sidebar.slider("Number of slides", 3, 12, 6)
max_bullets = st.sidebar.slider("Bullets per slide", 2, 6, 3)
temperature = st.sidebar.slider("Creativity (temperature)", 0.0, 1.0, 0.2, 0.05)

st.sidebar.markdown("---")
st.sidebar.write("If you select **Gemini**, paste your API key below.")
gemini_api_key = st.sidebar.text_input("Gemini API Key (optional)", type="password")

# ---------------------------- Prompt Helper --------------------------------
def build_prompt(topic: str, n_slides: int, bullets_per_slide: int) -> str:
    return (
        f"Create {n_slides} PowerPoint slides on the topic: \"{topic}\".\n"
        "Output ONLY a JSON list (no extra text) with this format:\n"
        "[\n"
        '  {"title":"Slide Title","bullets":["short bullet 1","short bullet 2"]},\n'
        "  ...\n"
        "]\n"
        f"Each slide should have {bullets_per_slide} bullets (short, <=12 words)."
    )

# ---------------------------- PPTX Builder ---------------------------------
def create_pptx(slides: List[Dict[str, Any]], deck_title: str) -> BytesIO:
    prs = Presentation()
    # Title slide
    try:
        title_slide_layout = prs.slide_layouts[0]
    except Exception:
        title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = deck_title

    # Content slides
    for s in slides:
        try:
            content_layout = prs.slide_layouts[1]
        except Exception:
            content_layout = prs.slide_layouts[0]
        sl = prs.slides.add_slide(content_layout)
        if sl.shapes.title:
            sl.shapes.title.text = s.get("title", "Untitled")
        # placeholder[1] is usually the body placeholder
        try:
            tf = sl.shapes.placeholders[1].text_frame
        except Exception:
            # fallback: create a textbox
            left = Inches(0.5) if 'Inches' in globals() else 0
            tf = sl.shapes.add_textbox(0, 0, 1, 1).text_frame  # improbable fallback

        tf.clear()
        bullets = s.get("bullets", [])
        for i, b in enumerate(bullets):
            if i == 0:
                tf.text = b
            else:
                p = tf.add_paragraph()
                p.text = b

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

# ---------------------------- HuggingFace Helpers --------------------------
@st.cache_resource(show_spinner=False)
def load_hf_pipeline(model_id: str):
    # delayed import to avoid heavy startup when not used
    from transformers import pipeline, AutoTokenizer, AutoModelForSeq2SeqLM
    # Load tokenizer and model, create text2text generation pipeline
    # Using CPU; Streamlit Cloud typically doesn't have GPU
    tokenizer = AutoTokenizer.from_pretrained(model_id)
    model = AutoModelForSeq2SeqLM.from_pretrained(model_id)
    gen = pipeline("text2text-generation", model=model, tokenizer=tokenizer, device_map="auto" if False else None)
    return gen

def generate_with_hf(model_id: str, prompt: str, max_new_tokens: int = 256, temperature: float = 0.2) -> str:
    gen = load_hf_pipeline(model_id)
    # The HF pipeline returns a dict with 'generated_text'
    out = gen(prompt, max_new_tokens=max_new_tokens, do_sample=temperature>0, temperature=temperature)
    if isinstance(out, list):
        return out[0].get("generated_text", "")
    return str(out)

# ---------------------------- Gemini Helper --------------------------------
def generate_with_gemini(api_key: str, prompt: str, model_name: str = "gemini-2.5-flash", temperature: float = 0.2):
    # Use google-generativeai if available
    try:
        import google.generativeai as genai
    except Exception as e:
        raise RuntimeError("google-generativeai is not installed in the environment.") from e

    genai.configure(api_key=api_key)
    model = genai.GenerativeModel(model_name)
    # This call may raise on quota/errors; bubble up to UI
    response = model.generate_content(prompt, temperature=temperature)
    # response.text contains the model output
    return response.text

# ---------------------------- Main Action ---------------------------------
topic = st.text_input("Enter your presentation topic", "")
generate_button = st.button("Generate Presentation 🚀")

if generate_button:
    if not topic.strip():
        st.error("Please enter a topic.")
    else:
        prompt = build_prompt(topic, num_slides, max_bullets)
        st.info("Prompt built. Generating slides...")
        st.write("**Prompt preview (hidden to model in production):**")
        with st.expander("Show prompt"):
            st.code(prompt)

        slides_json = None
        raw_output = None
        # Choose path
        try:
            if model_choice.startswith("HF"):
                # Choose model id
                hf_model = "google/flan-t5-large" if "Better" in model_choice else "google/flan-t5-base"
                st.info(f"Using HuggingFace model: {hf_model} (this runs on CPU; may take a few seconds)")
                # generate
                raw_output = generate_with_hf(hf_model, prompt, max_new_tokens=256, temperature=temperature)
            else:
                # Gemini selected
                if not gemini_api_key:
                    st.error("You chose Gemini but did not provide an API key in the sidebar.")
                    st.stop()
                st.info("Using Gemini (remote API). This may hit quota limits depending on your key.")
                raw_output = generate_with_gemini(gemini_api_key, prompt, model_name="gemini-2.5-flash", temperature=temperature)

            # Clean the raw output (strip markdown fences)
            if raw_output is None:
                raise RuntimeError("No output from model.")
            out = raw_output.strip()
            out = out.replace("```json", "").replace("```", "")
            # Attempt to parse JSON
            try:
                slides_json = json.loads(out)
            except Exception:
                # Try to find the first '[' ... ']' block
                import re
                m = re.search(r"(\[.*\])", out, re.S)
                if m:
                    candidate = m.group(1)
                    try:
                        slides_json = json.loads(candidate)
                    except Exception as e:
                        slides_json = None
                else:
                    slides_json = None

            if slides_json is None:
                st.warning("Model output wasn't valid JSON. We'll create fallback slides from the topic.")
                # create fallback slides
                slides_json = []
                for i in range(num_slides):
                    slides_json.append({
                        "title": f"{topic} - Slide {i+1}",
                        "bullets": [f"Auto bullet {j+1}" for j in range(max_bullets)]
                    })
                st.info("Fallback slides created.")
            else:
                st.success("Slides JSON parsed successfully.")

            # Show preview
            st.subheader("👀 Slide Preview")
            for i, s in enumerate(slides_json):
                st.markdown(f"**Slide {i+1} — {s.get('title','(no title)')}**")
                bullets = s.get("bullets", [])
                for b in bullets:
                    st.write("- " + b)

            # Create PPTX and offer download
            with st.spinner("Creating PPTX..."):
                pptx_bytes = create_pptx(slides_json, topic)
            filename = topic.replace(" ", "_")[:40] + ".pptx"
            st.download_button("📥 Download PPTX", data=pptx_bytes, file_name=filename,
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            st.success("Presentation ready!")

            # Optionally show raw model output (for debugging)
            with st.expander("Show raw model output"):
                st.code(raw_output)

        except Exception as e:
            st.error(f"An error occurred: {e}")
            # Show some help hints
            st.markdown("**Hints:**")
            st.markdown("- If using HF models, ensure `transformers` and `torch` are installed on the server.")
            st.markdown("- If using Gemini, ensure your key has quota and the model name is supported.")
            st.markdown("- For heavy models, running on CPU can be slow; try the HF Fast option for quicker results.")
