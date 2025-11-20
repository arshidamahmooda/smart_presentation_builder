from pptx import Presentation
from pptx.util import Pt
from io import BytesIO

icon = "🩺"

def generate_presentation_green_theme(slides):
    prs = Presentation()

    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = slides[0]["title"]

    # Content Slides
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{icon} {slide_data['title']}"

        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        for i, bullet in enumerate(slide_data["bullets"]):
            if i == 0:
                body.text = f"• {bullet}"
            else:
                p = body.add_paragraph()
                p.text = f"• {bullet}"
            p.font.size = Pt(22)

    ppt_data = BytesIO()
    prs.save(ppt_data)
    ppt_data.seek(0)
    return ppt_data
